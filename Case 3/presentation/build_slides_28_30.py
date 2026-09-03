"""
build_slides_28_30.py  -  three drop-in slides for All-Kanns_Asset_Management.pptx
(the master deck is hand-edited, so these are produced standalone in the same
visual language and copied in).

  28  Returns - what the 15-year Monte-Carlo produces           (chart + 2 boxes)
  29  Keeping within our risk budget - 1y 99% VaR vs the limit   (same layout)
  30  Splitting up the risk - standalone 99% loss by driver

Numbers pulled from  results_var/MC_LIFECYCLE_REPORT.md, HS_REPORT_full_(unhedged|
hedged).md, MC_REPORT_mc_full.md, SOLVENCY_II.md.  Charts from presentation/assets/.
Run:  python presentation/build_slides_28_30.py   ->  presentation/slides_28_30.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "slides_28_30.pptx"

NAVY   = RGBColor(0x12, 0x32, 0x3F)
DTEAL  = RGBColor(0x1F, 0x4E, 0x5F)
TEAL   = RGBColor(0x3B, 0x8E, 0x9E)
INK    = RGBColor(0x33, 0x33, 0x33)
MUTE   = RGBColor(0x8A, 0x8A, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PANEL  = RGBColor(0xED, 0xF2, 0xF4)
TEALL  = RGBColor(0xE7, 0xF0, 0xF2)
PINK   = RGBColor(0xE0, 0x80, 0xE0)
RULE   = RGBColor(0xC7, 0xD3, 0xD7)
FONT   = "Calibri"
FONT_L = "Calibri Light"

SECTIONS = ["Client Situation", "Our Selling Points", "Strategic Challenges & Market Overview",
            "Product Strategy", "Portfolio Results", "What Happens With the Upside"]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _run(p, text, size, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.05
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, s, kw) in para:
            _run(p, t, s, **kw)
    return tb


def bullets(slide, x, y, w, h, items, size=10, gap=5):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        _run(p, "•  ", size, color=TEAL)
        if isinstance(it, tuple):
            _run(p, it[0], size, color=DTEAL, bold=True)
            _run(p, it[1], size, color=INK)
        else:
            _run(p, it, size, color=INK)


def panel(slide, x, y, w, h, heading, items, size=10):
    box(slide, x, y, w, h, fill=PANEL, rounded=True)
    box(slide, x, y, w, 0.05, fill=TEAL)
    text(slide, x + 0.16, y + 0.13, w - 0.32, 0.4, [[(heading, 12.5, dict(color=DTEAL, bold=True))]])
    bullets(slide, x + 0.13, y + 0.63, w - 0.28, h - 0.75, items, size=size)


def stat_tiles(slide, y, tiles, x0=0.44, total_w=12.5, h=0.82):
    n = len(tiles)
    gap = 0.2
    w = (total_w - gap * (n - 1)) / n
    for i, (lab, big, sub) in enumerate(tiles):
        x = x0 + i * (w + gap)
        box(slide, x, y, w, h, fill=PANEL, rounded=True)
        box(slide, x, y, w, 0.045, fill=TEAL)
        text(slide, x + 0.13, y + 0.08, w - 0.24, 0.2, [[(lab.upper(), 7.5, dict(color=TEAL, bold=True))]])
        text(slide, x + 0.13, y + 0.28, w - 0.24, 0.34, [[(big, 15, dict(color=NAVY, bold=True))]])
        text(slide, x + 0.13, y + 0.61, w - 0.24, 0.2, [[(sub, 7.5, dict(color=INK))]], line_spacing=1.0)


def image(slide, name, x, y, w):
    return slide.shapes.add_picture(str(ASSETS / name), Inches(x), Inches(y), width=Inches(w))


def note(slide, txt, y=6.62, size=9):
    text(slide, 0.44, y, 12.5, 0.4, [[(txt, size, dict(color=MUTE, italic=True))]])


def chrome(title, subtitle, page, section_idx=4):
    s = prs.slides.add_slide(BLANK)
    box(s, 0.0, 0.30, 0.14, 0.62, fill=TEAL)
    box(s, 12.86, 0.20, 0.34, 0.34, fill=PINK)
    text(s, 0.42, 0.24, 11.9, 0.62, [[(title, 24, dict(color=NAVY, bold=True, font=FONT_L))]])
    text(s, 0.44, 0.82, 12.2, 0.40, [[(subtitle, 12.5, dict(color=TEAL))]])
    box(s, 0.44, 1.24, 12.5, 0.014, fill=RULE)
    text(s, 11.9, 0.20, 0.8, 0.3, [[(str(page), 12, dict(color=MUTE))]], align=PP_ALIGN.RIGHT)
    tb = s.shapes.add_textbox(Inches(0.44), Inches(7.04), Inches(12.4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    for i, sec in enumerate(SECTIONS):
        if i:
            _run(p, "   |   ", 8, color=RGBColor(0xB6, 0xC2, 0xC6))
        _run(p, sec.upper(), 8, color=(DTEAL if i == section_idx else RGBColor(0xB6, 0xC2, 0xC6)),
             bold=(i == section_idx))
    box(s, 0.0, 7.36, 13.333, 0.14, fill=DTEAL)
    return s


def banner(slide, label, y=1.44, w=12.5):
    box(slide, 0.44, y, w, 0.36, fill=DTEAL)
    text(slide, 0.60, y + 0.015, w - 0.3, 0.33, [[(label.upper(), 11, dict(color=WHITE, bold=True))]],
         anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================== #
#  28 - RETURNS FROM THE MONTE-CARLO
# ======================================================================== #
s = chrome("Returns – What the 15-Year Monte-Carlo Produces",
           "50,000 paths on the actual funding waterfall; returns above the 1.0% guarantee on every percentile", 28)
banner(s, "Monte-Carlo return distribution")
image(s, "lc_03_15year_IRR_distribution.png", 0.44, 1.98, 7.55)
panel(s, 8.20, 1.98, 4.74, 2.28, "How the returns are simulated", [
    ("50,000 paths, annual steps.  ", "€5.0bn into the matched bond book at t=0; €0.5bn/yr into the 14-index return portfolio, years 1–10."),
    ("Multivariate Student-t (df 5), ", "scaled so the shock covariance equals the 10-year historical covariance; geometric-mean drift; −99% floor."),
    ("Realised in the sim:  ", "return portfolio 11.1% mean / 16.5% vol p.a., bond book 3.9% / 5.2% — in line with history."),
], size=9.3)
panel(s, 8.20, 4.42, 4.74, 2.28, "What the returns say", [
    ("Median 15-year IRR 6.5% p.a.  ", "≈ 550 bps above the 1.0% guarantee; mean annual return 6.6%."),
    ("5th-percentile path 2.6% p.a., 0.5th-percentile 0.2% ", "— still positive, still ahead of the guaranteed accrual."),
    ("Median assets at year 15 ≈ €22.0bn ", "on €10.0bn of paid-in capital → median funded ratio 220%."),
], size=9.3)
stat_tiles(s, 5.78, [
    ("Median 15y IRR", "6.5% p.a.", "money-weighted"),
    ("5th-pct IRR", "2.6% p.a.", "clears the guarantee"),
    ("0.5th-pct IRR", "0.2% p.a.", "still ≥ 0"),
    ("Mean annual return", "6.6%", "contribution-stripped"),
], x0=0.44, total_w=7.55)
note(s, "Source: mc_lifecycle.py / results_var/MC_LIFECYCLE_REPORT.md.", y=6.9, size=8)

# ======================================================================== #
#  29 - KEEPING WITHIN THE RISK BUDGET
# ======================================================================== #
s = chrome("Keeping Within Our Risk Budget",
           "The 1-year 99% VaR against the board's economic risk limit and the Solvency II requirement", 29)
banner(s, "1-year 99% VaR vs. the risk limit")
image(s, "var_bridge.png", 0.44, 2.05, 7.55)
panel(s, 8.20, 1.98, 4.74, 2.28, "Where the limit comes from", [
    ("Own funds (economic surplus) €3.19bn.  ", "Board funding-ratio floor 1.20 → max tolerable 1-year asset loss €1.83bn."),
    ("Net of non-equity surplus risk, ", "the equity 99% 1y VaR limit is €0.83bn — rule-derived, not an arbitrary %."),
    ("Regulatory:  ", "SCR (99.5%, internal model) €0.94bn → Solvency II ratio ≈ 340%; MCR ≈ €0.23bn (> 1,300% covered)."),
], size=9.3)
panel(s, 8.20, 4.42, 4.74, 2.28, "Where we land", [
    ("Surplus 1y 99% VaR €1.21bn unhedged → €0.85bn ", "with the receiver-IRS overlay (−30%); Asset VaR €2.51bn."),
    ("Equity-futures overlay stays at 0%  ", "— the IRS already brings surplus risk inside the board limit."),
    ("Monte-Carlo (Student-t + Heston VIX) surplus VaR €1.81bn ", "— fatter tails, still well inside own funds."),
], size=9.3)
stat_tiles(s, 5.78, [
    ("Surplus VaR (hedged)", "€0.85bn", "−30% vs unhedged"),
    ("SCR (99.5%, 1y)", "€0.94bn", "internal model"),
    ("Solvency II ratio", "≈ 340%", "own funds / SCR"),
    ("Equity futures hedge", "0%", "IRS frees the budget"),
], x0=0.44, total_w=7.55)
note(s, "Source: results_var/HS_REPORT_full_(unhedged|hedged).md, HS_REPORT_scr_irs.md, MC_REPORT_mc_full.md, SOLVENCY_II.md.", y=6.9, size=8)

# ======================================================================== #
#  30 - SPLITTING UP THE RISK
# ======================================================================== #
s = chrome("Splitting Up the Risk",
           "Standalone 99% 1-year loss by driver — and why they net to a much smaller total", 30)
banner(s, "Standalone 99% 1-year loss by driver")
image(s, "risk_drivers.png", 0.44, 2.0, 7.7)
panel(s, 8.30, 1.98, 4.64, 2.28, "Rate risk nets out", [
    ("The guaranteed liability PV moves €2.92bn ", "and the bond book €1.54bn on the same rate shocks — in the surplus they largely cancel (Surplus P&L = ΔAssets − ΔLiability)."),
    ("The receiver IRS (€0.98bn standalone) ", "sits between them and closes the residual duration gap, so the surplus's net rate exposure is ≈ 0."),
], size=9.3)
panel(s, 8.30, 4.42, 4.64, 2.28, "The risks we choose to run", [
    ("Equity €0.96bn ", "— the deliberate risk budget, funded by surplus, never touching the guarantee."),
    ("FX-hedge residual €0.19bn ", "— the 1-year rate differential on rolling the USD→EUR hedges; HY and rates/credit indices negligible."),
    ("Longevity is not a market-VaR driver ", "— it is handled by buffer sizing and optional reinsurance."),
], size=9.3)
note(s, "The standalone losses sum to > €6bn but the diversified total surplus VaR is €0.85bn — the gap is diversification: "
        "rate moves hit assets and the liability together, and the IRS closes what is left.", y=6.66, size=9)

prs.save(str(OUT))
print("wrote", OUT.name, "|", len(prs.slides._sldIdLst), "slides")
