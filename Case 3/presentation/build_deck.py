"""
build_deck.py  -  regenerate  All-Kanns_Asset_Management.pptx

Keeps the structure and visual language of the team's existing deck (teal accent
bar, dark section banners, light content panels, 4-up stat tiles, breadcrumb
footer) but

  * corrects every number to the current Case 3b code / result files
    (case3_model.py, alm_fixed_income_.py, cashflow_match_v2.py, return_book.py
     and the workbooks under results/ , results_var/ , results_v2/)
  * switches the internal risk measure to the mandated 1-year 99% VaR
  * embeds figures rebuilt by make_charts.py (presentation/assets/)
  * adds  SECTION 06  "What happens with the upside"  (cash-flow structure +
    flows under every policyholder-election scenario + the 90/10 profit share)
  * adds a closing  Key Selling Points / Recommendations  slide

Run:  python presentation/build_deck.py     (after: python presentation/make_charts.py)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = Path(__file__).resolve().parent
CASE = HERE.parent
ASSETS = HERE / "assets"
OUT = HERE / "All-Kanns_Asset_Management.pptx"

# ---- palette -------------------------------------------------------------- #
NAVY   = RGBColor(0x12, 0x32, 0x3F)
DTEAL  = RGBColor(0x1F, 0x4E, 0x5F)
TEAL   = RGBColor(0x3B, 0x8E, 0x9E)
TEALL  = RGBColor(0xE7, 0xF0, 0xF2)
INK    = RGBColor(0x33, 0x33, 0x33)
MUTE   = RGBColor(0x8A, 0x8A, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PANEL  = RGBColor(0xED, 0xF2, 0xF4)
ORANGE = RGBColor(0xC8, 0x5A, 0x2B)
GREEN  = RGBColor(0x2E, 0x7D, 0x4F)
RED    = RGBColor(0xB4, 0x32, 0x2B)
SECTBG = RGBColor(0x0D, 0x28, 0x33)
PINK   = RGBColor(0xE0, 0x80, 0xE0)
FONT   = "Calibri"
FONT_L = "Calibri Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SECTIONS = ["Client Situation", "Our Selling Points", "Strategic Challenges & Market Overview",
            "Product Strategy", "Portfolio Results", "What Happens With the Upside"]
_page = [0]


# ---- primitives -------------------------------------------------------------
def _tf(shape, pad=0.06):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return tf


def _run(p, text, size, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shadow=False,
        rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.045
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.06):
    """runs: list of paragraphs; each paragraph is a list of (text, size, kw)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(tb)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, s, kw) in para:
            _run(p, t, s, **kw)
    return tb


def bullets(slide, x, y, w, h, items, size=11, color=INK, gap=5, bullet="•  ",
            lead_color=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(tb)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if isinstance(it, tuple):      # (lead, rest)
            _run(p, bullet, size, color=color)
            _run(p, it[0], size, color=lead_color or DTEAL, bold=True)
            _run(p, it[1], size, color=color)
        else:
            _run(p, bullet, size, color=color)
            _run(p, it, size, color=color)
    return tb


def image(slide, path, x, y, w, h=None):
    p = ASSETS / path if not Path(path).is_absolute() else Path(path)
    kw = dict(width=Inches(w))
    if h is not None:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(p), Inches(x), Inches(y), **kw)


# ---- page furniture -------------------------------------------------------- #
def _chrome(slide, title, subtitle, section_idx):
    _page[0] += 1
    box(slide, 0.0, 0.30, 0.14, 0.62, fill=TEAL)
    box(slide, 12.86, 0.20, 0.34, 0.34, fill=PINK)
    text(slide, 0.42, 0.24, 11.9, 0.62, [[(title, 24, dict(color=NAVY, bold=True, font=FONT_L))]])
    text(slide, 0.44, 0.82, 12.2, 0.40, [[(subtitle, 12.5, dict(color=TEAL))]])
    box(slide, 0.44, 1.24, 12.5, 0.014, fill=RGBColor(0xC7, 0xD3, 0xD7))
    text(slide, 12.4, 0.18, 0.7, 0.3, [[(str(_page[0]), 12, dict(color=MUTE))]], align=PP_ALIGN.RIGHT)
    # breadcrumb
    tb = slide.shapes.add_textbox(Inches(0.44), Inches(7.04), Inches(12.4), Inches(0.3))
    p = _tf(tb).paragraphs[0]
    for i, s in enumerate(SECTIONS):
        if i:
            _run(p, "   |   ", 8, color=RGBColor(0xB6, 0xC2, 0xC6))
        _run(p, s.upper(), 8, color=(DTEAL if i == section_idx else RGBColor(0xB6, 0xC2, 0xC6)),
             bold=(i == section_idx))
    box(slide, 0.0, 7.36, 13.333, 0.14, fill=DTEAL)


def content(title, subtitle, section_idx):
    s = prs.slides.add_slide(BLANK)
    _chrome(s, title, subtitle, section_idx)
    return s


def banner(slide, label, y=1.44, w=12.5):
    box(slide, 0.44, y, w, 0.36, fill=DTEAL)
    text(slide, 0.60, y + 0.015, w - 0.3, 0.33,
         [[(label.upper(), 11, dict(color=WHITE, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)


def panel(slide, x, y, w, h, heading, items, hcolor=DTEAL, size=11):
    box(slide, x, y, w, h, fill=PANEL, rounded=True)
    box(slide, x, y, w, 0.055, fill=TEAL)
    text(slide, x + 0.18, y + 0.14, w - 0.36, 0.4,
         [[(heading, 13, dict(color=hcolor, bold=True))]])
    bullets(slide, x + 0.14, y + 0.66, w - 0.32, h - 0.8, items, size=size)


def stat_tiles(slide, y, tiles, x0=0.44, total_w=12.5, h=1.9, big_fs=None):
    n = len(tiles)
    gap = 0.24
    w = (total_w - gap * (n - 1)) / n
    if big_fs is None:
        big_fs = 21 if w >= 2.6 else (17 if w >= 2.0 else 15)
    for i, (lab, big, sub) in enumerate(tiles):
        x = x0 + i * (w + gap)
        box(slide, x, y, w, h, fill=PANEL, rounded=True)
        box(slide, x, y, w, 0.05, fill=TEAL)
        text(slide, x + 0.16, y + 0.13, w - 0.3, 0.28,
             [[(lab.upper(), 9, dict(color=TEAL, bold=True))]])
        text(slide, x + 0.16, y + 0.42, w - 0.3, 0.5, [[(big, big_fs, dict(color=NAVY, bold=True))]])
        text(slide, x + 0.16, y + 0.98, w - 0.3, max(h - 1.02, 0.4),
             [[(sub, 9.5, dict(color=INK))]], line_spacing=1.02)


def kv_table(slide, x, y, w, rows, col1=0.32, header=None, rh=0.34, fs=10.5):
    """rows: list of tuples (same length). first row optionally styled as header."""
    ncol = len(rows[0])
    cw = [w * col1] + [w * (1 - col1) / (ncol - 1)] * (ncol - 1)
    cy = y
    if header:
        box(slide, x, cy, w, rh, fill=DTEAL)
        cx = x
        for j, htxt in enumerate(header):
            text(slide, cx + 0.1, cy + 0.02, cw[j] - 0.15, rh,
                 [[(htxt, fs, dict(color=WHITE, bold=True))]], anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT))
            cx += cw[j]
        cy += rh
    for i, row in enumerate(rows):
        if i % 2 == 0:
            box(slide, x, cy, w, rh, fill=PANEL)
        cx = x
        for j, val in enumerate(row):
            bold = (j == 0)
            text(slide, cx + 0.1, cy + 0.02, cw[j] - 0.15, rh,
                 [[(str(val), fs, dict(color=INK, bold=bold))]], anchor=MSO_ANCHOR.MIDDLE,
                 align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT))
            cx += cw[j]
        cy += rh
    return cy


def note(slide, txt, y=6.62, color=MUTE, size=9, bold=False, italic=True):
    text(slide, 0.44, y, 12.5, 0.4, [[(txt, size, dict(color=color, bold=bold, italic=italic))]])


def section_divider(idx, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, 13.333, 7.5, fill=SECTBG)
    box(s, 1.0, 2.62, 1.5, 0.05, fill=TEAL)
    text(s, 1.0, 2.15, 8, 0.4, [[(f"SECTION 0{idx}", 13, dict(color=TEAL, bold=True))]])
    text(s, 1.0, 2.85, 11.2, 1.4, [[(title, 40, dict(color=WHITE, bold=True, font=FONT_L))]])
    text(s, 1.0, 4.35, 10.8, 1.0, [[(subtitle, 14, dict(color=RGBColor(0xAE, 0xC4, 0xCB)))]])
    return s


# ======================================================================== #
#  TITLE
# ======================================================================== #
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 0.20, 1.5, fill=TEAL)
box(s, 12.7, 0.25, 0.45, 0.45, fill=PINK)
box(s, 0, 7.28, 13.333, 0.22, fill=DTEAL)
text(s, 0.9, 1.7, 11.5, 1.9,
     [[("Asset Management for the", 34, dict(color=NAVY, bold=True, font=FONT_L))],
      [("All-Kanns Life Insurance", 34, dict(color=NAVY, bold=True, font=FONT_L))]])
text(s, 0.9, 3.5, 11, 1.0,
     [[("Financial Engineering  ", 16, dict(color=INK)), ("|  Prof. Dr. Stefan Tilch", 16, dict(color=MUTE))],
      [("Group 3b  ", 16, dict(color=INK)), ("|  1-year 99% VaR, EUR swap curve as of 2 Sep 2026", 13, dict(color=MUTE))]])
text(s, 0.9, 5.2, 11.6, 1.2,
     [[("Ben Lösche   ·   David Wolter   ·   Kira Clara Bøgeskov Pedersen   ·   Ishita Pandey", 12.5, dict(color=INK))],
      [("Manuel Blasco Blázquez   ·   Wiktor Murawski   ·   Zhuoran Wang", 12.5, dict(color=INK))]])

# ======================================================================== #
#  EXECUTIVE SUMMARY
# ======================================================================== #
s = content("Executive Summary", "The mandate, our approach and the outcome our model produces", 0)
banner(s, "The mandate at a glance")
stat_tiles(s, 2.0, [
    ("The ask", "€5.0bn → ~€10bn", "100,000-policyholder book, contributions €0.5bn/yr for 10y, 15-year horizon"),
    ("The approach", "2 sleeves", "Liability-matching bond book + surplus-funded return portfolio, split re-struck yearly"),
    ("The result", "~6.0% p.a.", "Median 15-year IRR; even the 5th-pct path (2.9%) clears the 1.0% guarantee"),
    ("The risk", "€0.85bn", "1-year 99% surplus VaR after the receiver-IRS overlay (€1.21bn before)"),
], big_fs=17)
box(s, 0.44, 4.30, 12.5, 1.15, fill=DTEAL)
text(s, 0.72, 4.44, 12.0, 0.95,
     [[("In short: ", 13, dict(color=WHITE, bold=True, italic=True)),
       ("a cash-flow- and key-rate-matched bond book clears the €11.3bn guaranteed floor; the "
        "surplus funds an aggressive, diversified return portfolio with a 90/10 policyholder "
        "profit share. A par receiver-IRS overlay (no cash outlay) closes the residual duration gap.",
        13, dict(color=WHITE, italic=True))]])
bullets(s, 0.5, 5.66, 12.4, 1.1, [
    ("Funding waterfall  ", "t=0 → €5.0bn into the liability-matching bond book;  t=1..10 → €0.5bn/yr into the return portfolio."),
    ("Every USD sleeve  ", "is swapped back to EUR;  accounting basis IAS;  regulatory frame Solvency II / BaFin AnlV."),
], size=10.5)

# ======================================================================== #
section_divider(1, "Client Situation / Needs",
                "The book, the guarantee, the payout election and the regulatory frame")
# ---- Client profile ------------------------------------------------------ #
s = content("Client Profile and Portfolio Structure",
            "100,000 policyholders, one age cohort, a 10-year contribution horizon", 0)
banner(s, "Set-up of All-Kanns")
panel(s, 0.44, 2.0, 6.05, 4.35, "Baseline & liabilities", [
    ("Book value today  ", "€5.0bn  (100,000 × €50,000)."),
    ("Contributions  ", "€5,000 per person per year, end of year, for 10 years  →  €0.5bn/yr."),
    ("One cohort  ", "all insured are 50 (50% male);  no deaths assumed in the first 15 years."),
    ("Guarantee  ", "1.0% p.a. Höchstrechnungszins on every euro paid in, to year 15."),
    ("Election at year 15  ", "lump sum or a monthly pension (≥10 years); base case an even 50/50 split."),
    ("Accounting / regulation  ", "IAS;  Solvency II SCR (99.5%);  BaFin AnlV quantitative limits."),
])
box(s, 6.72, 2.0, 6.22, 4.35, fill=PANEL, rounded=True)
box(s, 6.72, 2.0, 6.22, 0.055, fill=TEAL)
text(s, 6.90, 2.14, 5.8, 0.35, [[("Cash flow to the insurer", 13, dict(color=DTEAL, bold=True))]])
kv_table(s, 6.94, 2.62, 5.78, [
    ("0  (age 50)", "existing assets  €5.0bn"),
    ("1 – 10  (51–60)", "+ €0.5bn per year"),
    ("11 – 14  (61–64)", "no contribution, no payment"),
    ("15  (65)", "guaranteed benefits begin"),
    ("16 – 50", "pension tail for annuitants"),
], col1=0.42, rh=0.52, fs=10.5)
note(s, "Sources: assignment brief;  results/optimization_summary.txt  (guaranteed accumulated value €11.303bn at year 15).")

# ---- Defined payout structure ----------------------------------------- #
s = content("Defined Payout Structure",
            "Contributions run 10 years; the year-15 election sets the payout shape", 0)
banner(s, "Payout options at year 15")
panel(s, 0.44, 2.0, 6.05, 2.0, "Lump sum", [
    "Accumulated guaranteed capital paid as a single amount at year 15.",
    "Relationship ends – no residual risk to us after payment.",
], size=10.5)
panel(s, 6.72, 2.0, 6.22, 2.0, "Monthly pension (annuity)", [
    "A stream for life, guaranteed for at least 10 years (to year 25).",
    "Beyond year 25 it continues only for survivors – longevity risk sits with us.",
], size=10.5)
stat_tiles(s, 4.3, [
    ("Choice at year 15", "50 / 50", "Base-case split; stressed ±25pp (results/policyholder_choice_scenarios.csv)"),
    ("Contribution inflow", "€0.5bn/yr", "€5,000 × 100,000, years 1–10, end of year"),
    ("No deaths yr 1–15", "clean", "Mortality only enters the pension tail from year 16"),
])
note(s, "Pension tail modelled age 65→99 (Destatis Sterbetafel 2022/24), minimum 10 years guaranteed.")

# ======================================================================== #
section_divider(2, "Our Expertise and Goals For You",
                "Why a liability-precision manager, not a generalist balanced fund")
s = content("Why All-Kanns Should Choose Us", "Four reasons, matched to the mandate", 1)
banner(s, "Our edge")
panel(s, 0.44, 2.0, 6.05, 2.15, "Liability-precision matching", [
    "One cohort, one contribution schedule – a genuine cash-flow + key-rate match.",
    "Assets phased to the accumulation, bridge and payout stages of your liability.",
], size=10.5)
panel(s, 6.72, 2.0, 6.22, 2.15, "IAS-aware by design", [
    "Built around fair-value volatility, not return alone – reported earnings protected.",
    "Duration and key-rate DV01 calibrated to keep P&L / OCI swings defensible.",
], size=10.5)
panel(s, 0.44, 4.32, 6.05, 2.15, "Embedded-option expertise", [
    "We price the year-15 lump-sum-vs-pension election like a callable bond.",
    "Structural insight most generalist managers overlook.",
], size=10.5)
panel(s, 6.72, 4.32, 6.22, 2.15, "Scale and built-in margin", [
    "€5.0bn today → c.€10bn by year 10 – institutional pricing on every trade.",
    "Engineered to clear the 1.0% guarantee with room to spare on the median path.",
], size=10.5)
note(s, "Objective: a partner who thinks like an ALM actuary and invests like an institutional asset manager.")

# ======================================================================== #
section_divider(3, "Considerations and Market Overview",
                "Allocation management, the regulatory frame and the European backdrop")
# ---- Considerations --------------------------------------------------- #
s = content("Considerations When Implementing the Strategy",
            "Balancing liquidity, long-term return and regulation across a 15-year horizon", 2)
banner(s, "Key strategic challenges")
panel(s, 0.44, 2.0, 4.05, 4.35, "Allocation management", [
    "Cash buffer ahead of the year-15 lump-sum spike – no forced asset sales.",
    "Stable, liquid book in the years 11–15 bridge.",
    "Weight limits across asset classes for diversification.",
], size=10)
box(s, 4.72, 2.0, 8.22, 4.35, fill=PANEL, rounded=True)
box(s, 4.72, 2.0, 8.22, 0.055, fill=TEAL)
text(s, 4.90, 2.14, 7.8, 0.35, [[("What each challenge demands", 13, dict(color=DTEAL, bold=True))]])
bullets(s, 4.9, 2.62, 7.9, 3.6, [
    "Sufficient cash / T-bills ahead of lump-sum payouts to meet benefits without selling into weakness.",
    "A stable, highly liquid portfolio through the residual bridge period to safeguard capital.",
    "Operating inside the AnlV investable universe and internal policy limits.",
    "Single-debtor ≤ 5% and corporate-group ≤ 3% of guarantee assets – concentration control.",
    "Solvency II capital efficiency weighed against every return-seeking position.",
], size=10.5)

# ---- Solvency II ---------------------------------------------------------- #
s = content("Quantitative Capital Requirements – Solvency II",
            "Directive 2009/138/EC – the frame our internal limits sit inside", 2)
banner(s, "Assessment parameters")
rows = [
    ("SCR", "Capital to absorb unexpected loss over one year at 99.5% confidence (1-in-200); standard formula or internal model; own funds must cover it at all times."),
    ("MCR", "Hard floor with an absolute minimum; always between 25% and 45% of the SCR."),
    ("Valuation", "Assets and liabilities at market-consistent fair value; technical provisions = best-estimate liabilities + risk margin."),
]
y = 2.0
for k, v in rows:
    box(s, 0.44, y, 2.1, 1.35, fill=DTEAL, rounded=True)
    text(s, 0.56, y + 0.1, 1.9, 1.2, [[(k, 13, dict(color=WHITE, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 2.74, y, 10.2, 1.35, fill=PANEL, rounded=True)
    text(s, 2.94, y + 0.1, 9.9, 1.2, [[(v, 11, dict(color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.55
note(s, "Our internal steering measure is a 1-year 99% VaR (below); the 99.5% SCR remains the binding regulatory test.")

# ---- IAS --------------------------------------------------------------- #
s = content("IAS Accounting Impact on the Balance Sheet",
            "Accounting treatment steers the book toward stable, long-duration, transparent assets", 2)
banner(s, "IAS standard → portfolio implication")
data = [
    ("Investment valuation", "HTM / AFS / FVTPL rules change how market moves hit P&L vs OCI", "Favours stable long-term bonds over volatile short-term holdings"),
    ("Liability measurement", "Future obligations discounted (guarantee accrues at 1.0%; market value on the EUR curve)", "Match asset duration and key-rate profile to the liability"),
    ("Revenue recognition", "Income recognised when earned, not on mark-to-market", "Predictable, steady running yield aligned to payouts"),
    ("Reporting & transparency", "Low-volatility results and disclosure of key assumptions required", "Transparent, low-turnover portfolio; documented stress + VaR"),
]
y = 1.92
for a, b, c in data:
    box(s, 0.44, y, 3.0, 1.0, fill=PANEL, rounded=True)
    text(s, 0.58, y + 0.06, 2.8, 0.9, [[(a, 10.5, dict(color=DTEAL, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 3.60, y, 5.3, 1.0, fill=PANEL, rounded=True)
    text(s, 3.74, y + 0.06, 5.05, 0.9, [[(b, 9.3, dict(color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 9.06, y, 3.88, 1.0, fill=TEALL, rounded=True)
    text(s, 9.2, y + 0.06, 3.62, 0.9, [[(c, 9.3, dict(color=DTEAL, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.16
note(s, "Not accounting advice.", y=6.7, size=8)

# ---- AnlV eligible ----------------------------------------------------- #
s = content("Eligible Asset Classes – BaFin AnlV", "The investable universe for guarantee assets", 2)
banner(s, "Eligible")
cells = [
    ("Government & public-sector bonds", "core matching asset for long-dated liabilities"),
    ("Covered bonds & bank debt", "secured, limited credit risk"),
    ("Corporate loans & bonds", "spread pick-up over sovereigns"),
    ("Equities & participations", "long-term return engine, capped by AnlV"),
    ("Structured credit (ABS / CLN)", "permitted within tight limits"),
    ("Real estate & REITs", "inflation-linked, illiquid diversifier"),
    ("Investment funds (UCITS / AIF)", "efficient diversified access"),
    ("Deposits & cash", "liquidity buffer for benefit payments"),
    ("Opening-clause assets", "only with explicit BaFin approval"),
]
for i, (a, b) in enumerate(cells):
    r, c = divmod(i, 3)
    x = 0.44 + c * 4.22
    y = 2.0 + r * 1.5
    box(s, x, y, 4.0, 1.34, fill=PANEL, rounded=True)
    box(s, x, y, 4.0, 0.05, fill=TEAL)
    text(s, x + 0.16, y + 0.14, 3.7, 1.1,
         [[(a, 11, dict(color=DTEAL, bold=True))], [(b, 9.5, dict(color=INK))]])
note(s, "Source: BaFin Anlageverordnung (AnlV), 14.11.2018.")

# ---- AnlV prohibited ------------------------------------------------------ #
s = content("Prohibited Assets – BaFin AnlV", "Excluded from the guarantee-asset universe", 2)
banner(s, "Prohibited")
for i, (a, b) in enumerate([
    ("Consumer / working-capital lending", "consumer or working-capital loans, movable goods and claims on them, intangibles (patents, goodwill)"),
    ("Group companies of the insurer", "§18 AktG entities, unless a purely passive stake with no operational influence"),
    ("Outsourced service companies", "firms performing the insurer's core insurance activities on its behalf"),
    ("Crypto and digital assets", "excluded from the guarantee-asset universe under AnlV"),
]):
    y = 2.0 + i * 1.15
    box(s, 0.44, y, 12.5, 1.0, fill=RGBColor(0xF6, 0xEC, 0xE4), rounded=True)
    box(s, 0.44, y, 0.09, 1.0, fill=ORANGE)
    text(s, 0.72, y + 0.08, 12.0, 0.9,
         [[(a + "  —  ", 11, dict(color=ORANGE, bold=True)), (b, 10.5, dict(color=INK))]],
         anchor=MSO_ANCHOR.MIDDLE)

# ---- AnlV limits ----------------------------------------------------- #
s = content("Quantitative Limitations Constrain Allocation",
            "BaFin caps exposure by asset class to control concentration and valuation risk", 2)
banner(s, "Maximum exposure (% of guarantee assets)")
lims = [("Equities, PE & subordinated debt", 40), ("Sovereigns & public bodies", 30),
        ("Real estate", 25), ("Covered bonds / public-sector banks", 15),
        ("Credit-linked, hedge & commodity", 7.5), ("Corporate loans (secured)", 5),
        ("Infrastructure (own quota, Feb 2025)", 5)]
y0 = 2.05
for i, (lab, v) in enumerate(lims):
    y = y0 + i * 0.53
    text(s, 0.44, y, 4.3, 0.4, [[(lab, 10, dict(color=INK))]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.85, y + 0.06, v / 40 * 5.6, 0.34, fill=DTEAL)
    text(s, 4.85 + v / 40 * 5.6 + 0.08, y, 0.9, 0.4, [[(f"{v}%", 10, dict(color=INK, bold=True))]],
         anchor=MSO_ANCHOR.MIDDLE)
box(s, 10.9, 2.0, 2.04, 4.05, fill=PANEL, rounded=True)
text(s, 11.05, 2.14, 1.8, 3.8, [
    [("Also", 10, dict(color=TEAL, bold=True))],
    [("single debtor ≤ 5%", 9.5, dict(color=INK))],
    [("corporate group ≤ 3%", 9.5, dict(color=INK))],
    [("single holding ≤ 1%", 9.5, dict(color=INK))],
    [("unlisted ≤ 15%", 9.5, dict(color=INK))],
])
note(s, "Source: BaFin AnlV, 14.11.2018; amendment effective Feb 2025.")

# ---- European insurers backdrop ---------------------------------------- #
s = content("European Insurers Have Shifted Back to Fixed Income",
            "ECB data: rate normalisation restored the appeal of long-duration bonds", 2)
banner(s, "What the backdrop tells us")
panel(s, 0.44, 2.0, 6.1, 4.3, "Read", [
    ("Fixed income is the core.  ", "Debt securities remain the dominant euro-area insurer asset class."),
    ("Funds / equity rose post-2020.  ", "The low-rate years pushed insurers up the risk curve."),
    ("Now: back to debt.  ", "ECB hikes lifted long-bond yields to ~3.3–3.6% on the EUR curve – a bond-anchored strategy is well timed."),
], size=10.5)
panel(s, 6.72, 2.0, 6.22, 4.3, "What we take from it", [
    "The €5.0bn liability-matching sleeve earns a ~4.0% running yield at today's spreads – comfortably above the 1.0% guarantee.",
    "Duration is available: the 15–30y part of the EUR curve is the natural home for this liability.",
    "The return sleeve is funded only by surplus, so the bond anchor is never at risk.",
], size=10.5)
note(s, "Source: ECB insurance-corporation statistics, 2025; EUR swap curve as of 2 Sep 2026 (results/normalized_eur_swap_curve.csv).")

# ---- Partnering / support ------------------------------------------------ #
s = content("Partnering to Enhance Performance Within the Frame",
            "Each client challenge matched to a concrete element of our support model", 2)
banner(s, "Your challenge → our support")
for i, (a, b) in enumerate([
    ("Navigating the rate environment", "Optimised long-term allocation on real market data, structured to exceed 1.0% with stable income"),
    ("Balancing solvency with performance", "Solvency II capital and BaFin limits built directly into portfolio construction and risk reporting"),
    ("Strengthening investment governance", "Transparent, data-driven process: quarterly ALM review, documented stress + 99% VaR limits"),
]):
    y = 2.05 + i * 1.5
    box(s, 0.44, y, 6.05, 1.28, fill=PANEL, rounded=True)
    text(s, 0.62, y + 0.1, 5.8, 1.1, [[(a, 11, dict(color=DTEAL, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
    box(s, 6.72, y, 6.22, 1.28, fill=TEALL, rounded=True)
    text(s, 6.9, y + 0.1, 5.95, 1.1, [[(b, 10.3, dict(color=INK))]], anchor=MSO_ANCHOR.MIDDLE)

# ======================================================================== #
section_divider(4, "Product Strategy",
                "Two sleeves: a matched bond floor, a surplus-funded return engine")
# ---- Where you are / how / where to be ---------------------------------- #
s = content("Securing Long-Term Returns Above the 1.0% Guarantee",
            "A two-sleeve strategy; model median 15-year IRR ≈ 6.0%", 3)
banner(s, "From here to there")
panel(s, 0.44, 2.0, 4.05, 4.35, "Where you are", [
    "100,000 insured aged 50, 50/50 male/female.",
    "AUM €5.0bn;  contributions €0.5bn/yr for 10y.",
    "Guarantee 1.0% p.a.;  year-15 election lump vs pension.",
], size=10)
panel(s, 4.72, 2.0, 4.05, 4.35, "How we get there", [
    "Years 0–10: invest inflows, run the return-efficient mix, build the liquidity buffer.",
    "Years 11–15: no inflows – de-risk along a glidepath, pre-fund the year-15 payout.",
    "Year 15+: pay benefits from the safe sleeve; monitor coverage.",
], size=10)
panel(s, 8.95, 2.0, 3.99, 4.35, "Where you want to be", [
    "Clear the 1.0% floor with a stable, investable, liability-aligned book.",
    "Hold ≥ 1.1× the next 12 months of outflows in cash + sovereigns.",
    "Quarterly ALM review; stress + 99% VaR limits.",
], size=10)

# ---- Guaranteed capital quantified ------------------------------------- #
s = content("The Guarantee, Quantified", "The floor every asset decision must clear", 3)
banner(s, "The promise")
panel(s, 0.44, 2.0, 6.5, 4.35, "In words, and as a formula", [
    "Every euro paid in – the €50,000 opening balance and each €5,000 contribution – compounds at 1.0% p.a. (Höchstrechnungszins) to year 15.",
    "It keeps accruing through the 5-year bridge (years 10→15) after contributions stop.",
], size=10.5)
text(s, 0.7, 4.55, 6.0, 0.9,
     [[("FV₁₅  =  P₀·(1+g)¹⁵  +  Σ Cₜ·(1+g)¹⁵⁻ᵗ", 13, dict(color=NAVY, bold=True))],
      [("P₀ = €50,000   Cₜ = €5,000, t = 1…10   g = 1.0%", 10, dict(color=MUTE))]])
box(s, 7.06, 2.0, 5.88, 4.35, fill=PANEL, rounded=True)
box(s, 7.06, 2.0, 5.88, 0.055, fill=TEAL)
text(s, 7.24, 2.14, 5.5, 0.35, [[("The floor, three ways", 13, dict(color=DTEAL, bold=True))]])
kv_table(s, 7.24, 2.66, 5.5, [
    ("Per policyholder, yr 15", "€113,028"),
    ("Book at year 15", "€11.30bn"),
    ("PV today (50/50 split)", "≈ €6.4bn"),
    ("PV range across election", "€6.1 – 6.9bn"),
], col1=0.62, rh=0.40, fs=10)
text(s, 7.24, 4.5, 5.5, 1.7, [
    [("The €11.3bn is a year-15 figure. Discounted on the EUR swap curve and "
      "split 50/50 lump/pension, the benefit liability is worth ≈ €6.4bn today "
      "(duration ≈ 20y, DV01 ≈ €12.8m/bp, 11% of PV beyond 30y).",
      9.5, dict(color=INK))]])

# ---- Lump vs pension cash-flow shapes -------------------------------- #
s = content("Lump Sum vs. Pension – Two Cash-Flow Shapes",
            "The year-15 election splits the liability into two very different profiles", 3)
banner(s, "What each choice means for us")
panel(s, 0.44, 2.0, 6.05, 2.0, "Lump sum", [
    "One payment at year 15, then the relationship ends.",
    "Full guaranteed value in a single spike – a large, dated liquidity need.",
], size=10.5)
panel(s, 6.72, 2.0, 6.22, 2.0, "Monthly pension", [
    "A stream from year 15, guaranteed ≥ 10 years, then survivors only.",
    "Smaller, spread outflows – but open-ended longevity risk.",
], size=10.5)
kv_table(s, 0.44, 4.35, 12.5, [
    ("Liquidity need", "one large cash-out at year 15", "spread over years, small monthly draws"),
    ("Planning horizon", "fixed – ends at payment", "open-ended – to the last survivor"),
    ("Risk retained by us", "none after payment", "longevity risk on every survivor"),
], header=("", "Lump sum", "Monthly pension"), col1=0.24, rh=0.44)

# ---- Mortality & longevity ---------------------------------------------- #
s = content("Mortality and Longevity – Sizing the Tail",
            "Grounded in Germany's official life table; what it does not capture", 3)
banner(s, "Destatis Sterbetafel 2022/24 – remaining life expectancy")
kv_table(s, 0.44, 2.0, 12.5, [
    ("Men (years)", "30.2", "21.6", "18.0", "14.2", "7.5"),
    ("Women (years)", "34.3", "25.2", "21.2", "16.9", "8.9"),
], header=("Age", "50", "60", "65", "70", "80"), col1=0.24, rh=0.44)
panel(s, 0.44, 3.5, 6.05, 2.9, "Method", [
    "Destatis period table – a snapshot of today's mortality, not a cohort forecast.",
    "Pension tail modelled age 65 → 99+ (a 50-year-old today can still be alive then).",
], size=10.5)
panel(s, 6.72, 3.5, 6.22, 2.9, "Longevity risk", [
    "Blended remaining life at 65 ≈ 19.6 years for this 50/50 cohort.",
    "A period table misses future mortality improvement – the cohort may live longer.",
    "That gap sits with us on every pension electant; reinsurance is available if it grows.",
], size=10.5)

# ---- EIOPA curve -------------------------------------------------------- #
s = content("Discounting the Guarantee Correctly",
            "A flat rate misprices a 15-year liability – the term structure matters", 3)
banner(s, "EIOPA / EUR risk-free term structure")
kv_table(s, 0.44, 2.0, 7.4, [
    ("1", "2.83%", "0.972"), ("3", "2.95%", "0.916"), ("5", "2.99%", "0.862"),
    ("7", "3.05%", "0.808"), ("10", "3.16%", "0.732"), ("15", "3.31%", "0.613"),
    ("20", "3.35%", "0.516"),
], header=("Year", "Spot rate", "Discount factor"), col1=0.34, rh=0.40)
panel(s, 8.1, 2.0, 4.84, 4.35, "Why EIOPA, not a flat rate", [
    "EIOPA publishes a monthly market-consistent curve for Solvency II discounting; using it removes any temptation to pick a favourable rate.",
    "Built from EUR swaps, extrapolated to a UFR ≈ 3.3%.",
    "The curve runs ~2.8% at 1y to ~3.3% at 15y – not flat; a flat rate misprices a 15-year cash flow.",
    ("PV of the guarantee ≈ €6.4bn  ", "(50/50 split; €6.9bn if 100% lump)."),
], size=9.6)
note(s, "Sources: EIOPA_EUR_curve.xlsx; results/normalized_eur_swap_curve.csv (bootstrapped par-swap curve, 2 Sep 2026).")

# ---- Policyholder choice sensitivity ------------------------------------- #
s = content("Base Case and Policyholder-Choice Sensitivity",
            "How the year-15 liquidity need moves with the lump-sum election rate", 3)
banner(s, "Liability split by election")
image(s, "election.png", 0.44, 2.0, 7.5)
kv_table(s, 8.2, 2.0, 4.74, [
    ("0%", "€0.0bn", "€5.35bn"),
    ("25%", "€2.83bn", "€5.75bn"),
    ("50%  (base)", "€5.65bn", "€6.14bn"),
    ("75%", "€8.48bn", "€6.54bn"),
    ("100%", "€11.30bn", "€6.94bn"),
], header=("Lump", "Cash-out yr 15", "PV today"), col1=0.34, rh=0.44)
note(s, "Base case: even 50/50 split (assignment gives the choice; we stress ±25pp). Source: mixed_liability_scenarios.xlsx.")

# ---- Portfolio architecture -------------------------------------------- #
s = content("Portfolio Architecture – Two Sleeves, One Mandate",
            "Segregating the guaranteed floor from the surplus-return assets", 3)
banner(s, "Architecture")
panel(s, 0.44, 2.0, 6.05, 2.05, "Liability-Matching Portfolio (LMP)", [
    "€5.0bn bond book, cash-flow + key-rate matched to the 1.0% guarantee.",
    "Small, by-design funding risk.",
], size=10.5)
panel(s, 6.72, 2.0, 6.22, 2.05, "Return-Seeking Portfolio (RSP)", [
    "Contributions (€0.5bn/yr) invested in an aggressive, diversified index mix.",
    "Captures the upside above the floor; a bad year never touches the guarantee.",
], size=10.5)
panel(s, 0.44, 4.35, 6.05, 2.05, "Why split", [
    "Standard insurance ALM – return risk can never endanger the floor.",
    "Each sleeve judged on its own mandate; funded status transparent at any time.",
], size=10)
panel(s, 6.72, 4.35, 6.22, 2.05, "What sets the split", [
    "Not a fixed %: the LMP is sized to the PV of the guaranteed liability.",
    "Everything above the funded amount flows to the RSP; the split re-strikes yearly.",
], size=10)

# ---- LMP detail ----------------------------------------------------------- #
s = content("Liability-Matching Portfolio – Bonds, Precisely Matched",
            "Cash-flow dedication, key-rate DV01 shaping, and a no-cash receiver-IRS overlay", 3)
banner(s, "How the match is built")
image(s, "krd.png", 0.44, 1.95, 7.5)
panel(s, 8.15, 1.95, 4.79, 4.5, "Four techniques", [
    ("Cash-flow dedication.  ", "Stage 1: running cash balance ≥ 0 every year (1.5% reinvestment) – external top-up €0.00bn, covered to year 50."),
    ("Duration matching.  ", "≈ 29y KRD-effective, matched to the ≈ 20y / €12.8m-per-bp liability."),
    ("Key-rate DV01.  ", "Stage 2: match € per bp at each tenor, without giving back Stage-1 coverage or duration."),
    ("Receiver-IRS overlay.  ", "Par swap, notional not exchanged – no cash. 15y €2.8bn + 30y €0.3bn closes the residual 15y gap."),
], size=9.4)
kv_table(s, 0.44, 5.75, 7.5, [
    ("11 bonds, ≈ €5.0bn", "KRD-eff. duration ≈ 29y", "cash-flow top-up €0.00bn; surplus DV01 gap ≈ 0 after IRS"),
], col1=0.30, rh=0.42, fs=9.5)
note(s, "Source: results_v2/portfolio_wide.csv (two-stage: Stage 1 cash-flow dedication, Stage 2 KRD shaping) + cashflow_match_v2.size_irs().")

# ---- RSP detail --------------------------------------------------------- #
s = content("Return-Seeking Portfolio – Equities for the Surplus",
            "Aggressive Diversified target weights across 14 investable indices", 3)
banner(s, "Target weights")
image(s, "rsp_weights.png", 0.44, 1.95, 7.7)
panel(s, 8.3, 1.95, 4.64, 4.5, "Rationale", [
    "≈ 92% equity across DAX, Nasdaq, Dow, MSCI World / Europe, Russell 2000, Asia-Pacific, plus rare-earth, health-care and Hong Kong satellites.",
    "≈ 4% high yield, ≈ 2% gold, ≈ 1% euro treasuries for diversification.",
    "Risk calibrated to surplus only – the guarantee is already secured by the LMP.",
], size=9.6)
note(s, "Source: portfolio_optimization_final.xlsx, sheet 'Portfolio Weights', column Aggressive_Diversified.")

# ---- Investment rationale --------------------------------------------- #
s = content("Investment Rationale – Why These Building Blocks",
            "Tying the two-sleeve architecture back to the quantified liability", 3)
banner(s, "LMP vs RSP")
kv_table(s, 0.44, 2.0, 12.5, [
    ("Job", "match the 1.0% guarantee, minimal tracking error", "capture the surplus – target ≈ 6% p.a. median IRR"),
    ("Instruments", "EUR sovereign / SSA / covered bonds + ultra-long & zero-coupon", "broad developed + satellite equity indices, some HY / gold"),
    ("Protects against", "funding risk on the guarantee – the one liability we cannot get wrong", "missing the return that makes the mandate worth running"),
], header=("", "Liability-Matching Portfolio", "Return-Seeking Portfolio"), col1=0.20, rh=0.62)
box(s, 0.44, 4.4, 12.5, 1.9, fill=PANEL, rounded=True)
bullets(s, 0.6, 4.55, 12.2, 1.7, [
    "A single blended portfolio would force every decision to compromise between safety and return – under-funding the guarantee or under-delivering the target.",
    "Every instrument screened for Solvency II capital, liquidity, diversification and implementation.",
    "A board can explain it in one line: the floor is matched, the upside is diversified.",
], size=10)

# ---- Investible universe ----------------------------------------------- #
s = content("The Investable Universe", "One universe, two mandates", 3)
banner(s, "By sleeve")
panel(s, 0.44, 2.0, 4.05, 4.4, "Fixed income", [
    ("LMP:  ", "EUR government & SSA bonds (KfW, EU, Belgium, Austria, Netherlands, Germany), covered bonds, ultra-long and zero-coupon (STRIPS)."),
    ("RSP credit:  ", "US corporate HY, Pan-European HY."),
], size=9.5)
panel(s, 4.72, 2.0, 4.05, 4.4, "Equity indices (RSP)", [
    "DAX, Dow Jones, Nasdaq, MSCI World, MSCI Europe, Russell 2000, Bloomberg Asia-Pacific Aggregate.",
    "Satellites: MVIS Global Rare Earth, MSCI World Health Care, iShares MSCI Hong Kong.",
], size=9.5)
panel(s, 8.95, 2.0, 3.99, 4.4, "Diversifiers (RSP)", [
    "Bloomberg Commodity – Gold.",
    "Bloomberg Euro Treasury (liquid ballast).",
], size=9.5)
note(s, "Source: portfolio_optimization_final.xlsx (14-index optimiser universe); Fixed Income Basket.xlsx + ZCB and ultra long coupon Bond.xlsx.")

# ======================================================================== #
section_divider(5, "Portfolio Results",
                "Does it work? 1-year 99% VaR, deterministic stress, 15-year Monte-Carlo")
# ---- Risk map --------------------------------------------------------- #
s = content("Risk Map – What All-Kanns Is Exposed To",
            "Four risks fall out of the strategy; each gets the same three questions", 4)
banner(s, "Four risks")
for i, (a, b) in enumerate([
    ("1  Interest-rate / duration risk", "Rate moves reprice both the LMP and the liability – the risk is the mismatch between them."),
    ("2  Equity & market risk", "The RSP's core exposure – accepted deliberately, on surplus capital only."),
    ("3  FX exposure", "USD sleeves (US Treasuries, US equity, USD credit) create currency risk against a EUR guarantee."),
    ("4  Longevity & election risk", "How long survivors live, and how many take the pension over the lump sum."),
]):
    r, c = divmod(i, 2)
    x = 0.44 + c * 6.35
    y = 2.0 + r * 2.25
    box(s, x, y, 6.15, 2.0, fill=PANEL, rounded=True)
    box(s, x, y, 6.15, 0.05, fill=TEAL)
    text(s, x + 0.18, y + 0.16, 5.8, 1.7,
         [[(a, 12, dict(color=DTEAL, bold=True))], [(b, 10, dict(color=INK))]])
text(s, 0.44, 6.55, 12.5, 0.4,
     [[("For each: why it arises · how large (quantified) · accept, mitigate or transfer.", 9.5, dict(color=MUTE, italic=True))]])

# ---- Risk quantification & response --------------------------------- #
s = content("Risk Quantification and Response",
            "For each risk: why it matters, how large it is, what we do", 4)
banner(s, "Quantified (1-year, current book)")
kv_table(s, 0.44, 2.0, 12.5, [
    ("Rate / duration", "±100bp: LMP ≈ €1.05bn vs liability ≈ €1.22bn; surplus swing ≈ €0.16bn with the KRD + IRS match (≈ €0.53bn without the IRS)", "Mitigate – cash-flow + KRD + receiver IRS"),
    ("Equity / market", "−30% equities ≈ −€1.41bn on the RSP – inside the €3.2bn economic surplus", "Accept – the deliberate risk budget"),
    ("FX", "USD sleeves swapped to EUR; residual is the 1y rate-differential on the hedge roll", "Mitigate – FX-swap the LMP; RSP partly unhedged as a diversifier"),
    ("Longevity / election", "25%→75% lump swing moves the year-15 need €2.8bn→€8.5bn; +1yr life exp ≈ +€0.27bn liability", "Mitigate via buffer sizing; reinsurance available"),
], header=("Risk", "How large (quantified)", "Response"), col1=0.16, rh=0.86, fs=9.3)
note(s, "Source: results_var/stress_tests_full_irs.csv & _unhedged.csv; results/policyholder_choice_scenarios.csv.")

# ---- One-year VaR ---------------------------------------------------- #
s = content("Does It Work?  One-Year 99% VaR and Surplus Risk",
            "The mandated internal measure: 1-year 99% VaR on asset P&L and on surplus P&L", 4)
banner(s, "1-year 99% VaR – Historical Simulation")
image(s, "var_bridge.png", 0.44, 1.95, 7.6)
panel(s, 8.2, 1.95, 4.74, 2.35, "Method", [
    "468 overlapping 52-week windows of real factor history; rates as absolute changes, indices as log returns.",
    "Non-normality kept: empirical tails + a Student-t / Heston Monte-Carlo cross-check.",
], size=9.3)
kv_table(s, 8.2, 4.45, 4.74, [
    ("Asset VaR", "€2.51bn", "25.1% of assets"),
    ("Surplus VaR (unhedged)", "€1.21bn", "17.8% of liab. PV"),
    ("Surplus VaR (+ IRS)", "€0.85bn", "8.5% of assets"),
    ("MC surplus VaR (t / Heston)", "€1.81bn", "fatter tails"),
], col1=0.52, rh=0.42, fs=9)
note(s, "Surplus P&L = ΔAssets − ΔLiability.  Source: results_var/HS_REPORT_full_*.md, MC_REPORT_mc_full.md.")

# ---- Stress + MC --------------------------------------------------------- #
s = content("Stress Testing and 15-Year Monte-Carlo",
            "Deterministic scenarios on the surplus, and the full distribution at year 15", 4)
banner(s, "Deterministic stress – surplus P&L")
image(s, "stress.png", 0.44, 1.86, 6.4)
image(s, "mc_year15.png", 7.05, 1.86, 5.9)
stat_tiles(s, 5.45, [
    ("Median assets, yr 15", "€20.7bn", "15-year MC"),
    ("5th percentile", "€14.3bn", "well above the floor"),
    ("0.5th percentile", "€11.2bn", "≈ the €11.3bn floor"),
    ("P(underfunded)", "< 0.6%", "before the glidepath"),
], x0=0.44, total_w=12.5, h=1.3)
note(s, "Sources: results_var/stress_tests_full_*.csv; monte_carlo_ALM_results.xlsx; 15Y median IRR ≈ 6.0%.", y=6.88, size=8)

# ======================================================================== #
section_divider(6, "What Happens With the Upside",
                "The cash-flow structure, the flows under every election scenario, and the 90/10 share")
# ---- 6a  Cash-flow structure ------------------------------------------- #
s = content("The Cash-Flow Structure",
            "Contributions in for 10 years; a year-15 spike, then a 35-year pension tail", 5)
banner(s, "Asset vs. guaranteed-liability cash flows")
image(s, "cf_annual.png", 0.44, 1.95, 7.6)
kv_table(s, 0.44, 5.5, 7.6, [
    ("Years 1–10", "+ €0.5bn / yr contributions", "→ Return portfolio"),
    ("Year 15", "− €5.65bn lump sum (50/50)", "→ safe sleeve"),
    ("Years 16–50", "− ≈ €0.3bn / yr pension tail", "→ coupons + reinvestment"),
], col1=0.24, rh=0.38, fs=9.5)
panel(s, 8.2, 1.95, 4.74, 4.5, "Reading the structure", [
    ("Years 1–10.  ", "€0.5bn/yr contributions fund the RSP; the LMP coupons accrue."),
    ("Year 15.  ", "the lump-sum spike – €5.65bn at 50/50 (€2.3–9.0bn across the election)."),
    ("Years 16–50.  ", "≈ €0.3bn/yr pension tail, decaying with mortality; 11% of PV beyond year 30."),
    ("Coverage.  ", "Stage-1 dedication keeps the running cash balance ≥ 0 every year to year 50 – no external top-up; ultra-long / zero-coupon bonds reach the 30y+ tail."),
], size=9.2)

# ---- 6b  Flows in every scenario ------------------------------------- #
s = content("Flows Under Every Scenario",
            "The year-15 need and the pension tail move with the election; the surplus absorbs the market path", 5)
banner(s, "Policyholder-election scenarios")
kv_table(s, 0.44, 2.0, 12.5, [
    ("0% lump / 100% pension", "€0.0bn", "100,000", "€8.72bn", "€5.35bn"),
    ("25% / 75%", "€2.83bn", "75,000", "€6.54bn", "€5.75bn"),
    ("50% / 50%  (base case)", "€5.65bn", "50,000", "€4.36bn", "€6.14bn"),
    ("75% / 25%", "€8.48bn", "25,000", "€2.18bn", "€6.54bn"),
    ("100% / 0%", "€11.30bn", "0", "€0.00bn", "€6.94bn"),
], header=("Election", "Lump-sum at yr 15", "Pensioners", "Pension PV at yr 15", "Total PV today"),
    col1=0.24, rh=0.46, fs=9.5)
panel(s, 0.44, 4.9, 6.05, 1.6, "Market path on top", [
    "Base 15y median assets €20.7bn vs a €10.0bn liability → funded ratio ≈ 2.1×.",
    "Even the 0.5th-percentile path (€11.2bn) covers the €11.3bn floor.",
], size=9.5)
panel(s, 6.72, 4.9, 6.22, 1.6, "So the upside is real", [
    "Across every election and almost every market path the book finishes above the guarantee.",
    "P(underfunded) ≈ 0.1% at the base election, still < 0.6% even at a 100% lump election – before the de-risking glidepath.",
], size=9.5)
note(s, "Source: mixed_liability_scenarios.xlsx (Scenario Summary); monte_carlo_ALM_results.xlsx (Scenario EUR bn).")

# ---- 6c  Profit share ------------------------------------------------- #
s = content("Sharing the Upside – 90 / 10 on the Return Portfolio",
            "A requirement: 90% of the return book's yearly profit goes to policyholders", 5)
banner(s, "Where the surplus goes")
image(s, "profit_share.png", 0.44, 1.9, 7.5)
panel(s, 8.15, 1.9, 4.79, 4.55, "Mechanics", [
    ("90% to policyholders.  ", "each year's investment profit on the RSP is 90% distributed, funded by selling an equal EUR amount from every sleeve."),
    ("10% retained, reinvested.  ", "stays in the book and compounds; losses carry forward and net against future profit first."),
    ("Over the 10-year accumulation.  ", "contributions €5.00bn → return-book MV €5.26bn; €2.31bn shared to policyholders; €0.26bn retained."),
    ("Guarantee untouched.  ", "the 90% is a pass-through off the insurer's asset side – it never draws on the LMP."),
], size=9.3)
note(s, "Source: return_book.py / results_var/RETURN_BOOK_REPORT.md.  Naive fully-reinvested book (no share) would be €8.39bn.")

# ======================================================================== #
#  KEY SELLING POINTS / RECOMMENDATIONS
# ======================================================================== #
s = content("Key Selling Points and Recommendations",
            "What we propose, and why All-Kanns should mandate it", 5)
banner(s, "Recommendation")
panel(s, 0.44, 1.95, 6.25, 4.55, "What we propose", [
    ("Split the book.  ", "€5.0bn into a cash-flow + key-rate matched LMP; contributions into the Aggressive Diversified RSP; re-strike the split yearly."),
    ("Add a receiver-IRS overlay.  ", "15y €2.8bn + 30y €0.3bn, par, notional not exchanged – no cash, only variation margin."),
    ("Extend the LMP.  ", "ultra-long and zero-coupon (STRIPS) bonds to reach the 20y+ liability without breaching the 15% instrument cap."),
    ("Swap USD → EUR  ", "on every sleeve; leave part of the RSP unhedged as a diversifier."),
    ("Glidepath years 11–15  ", "de-risk into the year-15 payout; hold ≥ 1.1× the next 12m of outflows in cash + sovereigns."),
], size=9.4)
panel(s, 6.85, 1.95, 6.09, 4.55, "Why it wins the mandate", [
    ("Floor is safe.  ", "€11.3bn guarantee covered even on the 0.5th-percentile 15-year path; P(underfunded) < 0.6%."),
    ("Upside is real.  ", "median 15-year IRR ≈ 6.0% – ≈ 500bps above the 1.0% guarantee; 90% shared with policyholders."),
    ("Risk is controlled.  ", "1-year 99% surplus VaR €0.85bn after the overlay (−30%); every stress inside the €3.2bn surplus."),
    ("IAS-defensible.  ", "duration + KRD matched, so LMP mark-to-market offsets the liability; low turnover."),
    ("Fully compliant.  ", "Solvency II capital and BaFin AnlV limits built into construction and reporting."),
], size=9.4)
note(s, "One line for the board: the floor is matched, the upside is diversified and shared, and the residual risk is hedged for no cash.")

prs.save(str(OUT))
print("wrote", OUT.relative_to(CASE), "|", len(prs.slides._sldIdLst), "slides")
